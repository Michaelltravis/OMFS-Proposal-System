"""
Google Drive API integration service
"""
import json
import io
from datetime import datetime
from typing import List, Optional, Dict, Any
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import Flow
from googleapiclient.discovery import build
from googleapiclient.errors import HttpError
from googleapiclient.http import MediaIoBaseDownload
from sqlalchemy.orm import Session
import pdfplumber
from docx import Document

from app.core.config import settings
from app.models.google_drive import GoogleDriveCredential
from app.schemas.google_drive import GoogleDriveFile


class GoogleDriveService:
    """Service for interacting with Google Drive API"""

    SCOPES = [
        "https://www.googleapis.com/auth/drive.readonly",
        "https://www.googleapis.com/auth/drive.metadata.readonly",
    ]

    @staticmethod
    def get_authorization_url(state: Optional[str] = None) -> str:
        """
        Generate Google OAuth authorization URL

        Args:
            state: Optional state parameter for CSRF protection

        Returns:
            Authorization URL for user to grant access
        """
        if not settings.GOOGLE_CLIENT_ID or not settings.GOOGLE_CLIENT_SECRET:
            raise ValueError("Google OAuth credentials not configured")

        flow = Flow.from_client_config(
            {
                "web": {
                    "client_id": settings.GOOGLE_CLIENT_ID,
                    "client_secret": settings.GOOGLE_CLIENT_SECRET,
                    "auth_uri": "https://accounts.google.com/o/oauth2/auth",
                    "token_uri": "https://oauth2.googleapis.com/token",
                    "redirect_uris": [settings.GOOGLE_REDIRECT_URI],
                }
            },
            scopes=GoogleDriveService.SCOPES,
        )
        flow.redirect_uri = settings.GOOGLE_REDIRECT_URI

        authorization_url, _ = flow.authorization_url(
            access_type="offline",
            include_granted_scopes="true",
            state=state,
            prompt="consent",  # Force consent to get refresh token
        )

        return authorization_url

    @staticmethod
    def exchange_code_for_token(code: str, db: Session) -> Dict[str, Any]:
        """
        Exchange authorization code for access token

        Args:
            code: Authorization code from OAuth callback
            db: Database session

        Returns:
            Token information and user details
        """
        if not settings.GOOGLE_CLIENT_ID or not settings.GOOGLE_CLIENT_SECRET:
            raise ValueError("Google OAuth credentials not configured")

        flow = Flow.from_client_config(
            {
                "web": {
                    "client_id": settings.GOOGLE_CLIENT_ID,
                    "client_secret": settings.GOOGLE_CLIENT_SECRET,
                    "auth_uri": "https://accounts.google.com/o/oauth2/auth",
                    "token_uri": "https://oauth2.googleapis.com/token",
                    "redirect_uris": [settings.GOOGLE_REDIRECT_URI],
                }
            },
            scopes=GoogleDriveService.SCOPES,
        )
        flow.redirect_uri = settings.GOOGLE_REDIRECT_URI

        # Exchange code for credentials
        flow.fetch_token(code=code)
        credentials = flow.credentials

        # Store credentials in database
        GoogleDriveService._save_credentials(db, credentials)

        # Get user info
        user_info = GoogleDriveService._get_user_info(credentials)

        return {
            "access_token": credentials.token,
            "expires_at": credentials.expiry,
            "user_email": user_info.get("emailAddress"),
        }

    @staticmethod
    def _save_credentials(db: Session, credentials: Credentials) -> None:
        """Save or update Google Drive credentials in database"""
        # Deactivate any existing credentials
        db.query(GoogleDriveCredential).update({"is_active": False})

        # Create new credential record
        credential = GoogleDriveCredential(
            access_token=credentials.token,
            refresh_token=credentials.refresh_token,
            token_uri=credentials.token_uri,
            client_id=credentials.client_id,
            client_secret=credentials.client_secret,
            scopes=json.dumps(credentials.scopes) if credentials.scopes else None,
            expiry=credentials.expiry,
            is_active=True,
        )
        db.add(credential)
        db.commit()

    @staticmethod
    def _get_credentials(db: Session) -> Optional[Credentials]:
        """Retrieve active Google Drive credentials from database"""
        credential = (
            db.query(GoogleDriveCredential)
            .filter(GoogleDriveCredential.is_active == True)
            .first()
        )

        if not credential:
            return None

        creds = Credentials(
            token=credential.access_token,
            refresh_token=credential.refresh_token,
            token_uri=credential.token_uri,
            client_id=credential.client_id,
            client_secret=credential.client_secret,
            scopes=json.loads(credential.scopes) if credential.scopes else None,
        )

        # Refresh if expired
        if creds.expired and creds.refresh_token:
            from google.auth.transport.requests import Request

            creds.refresh(Request())
            # Update database with new token
            credential.access_token = creds.token
            credential.expiry = creds.expiry
            db.commit()

        return creds

    @staticmethod
    def _get_user_info(credentials: Credentials) -> Dict[str, Any]:
        """Get user information from Google Drive API"""
        try:
            service = build("drive", "v3", credentials=credentials)
            about = service.about().get(fields="user").execute()
            return about.get("user", {})
        except HttpError:
            return {}

    @staticmethod
    def get_connection_status(db: Session) -> Dict[str, Any]:
        """
        Check if Google Drive is connected

        Returns:
            Connection status and user info
        """
        credential = (
            db.query(GoogleDriveCredential)
            .filter(GoogleDriveCredential.is_active == True)
            .first()
        )

        if not credential:
            return {"connected": False}

        creds = GoogleDriveService._get_credentials(db)
        if not creds:
            return {"connected": False}

        user_info = GoogleDriveService._get_user_info(creds)

        return {
            "connected": True,
            "user_email": user_info.get("emailAddress"),
            "expires_at": credential.expiry,
            "folder_id": credential.folder_id,
        }

    @staticmethod
    def disconnect(db: Session) -> None:
        """Disconnect Google Drive by deactivating credentials"""
        db.query(GoogleDriveCredential).update({"is_active": False})
        db.commit()

    @staticmethod
    def set_folder_id(db: Session, folder_id: Optional[str]) -> None:
        """
        Set the folder ID to restrict searches

        Args:
            db: Database session
            folder_id: Google Drive folder ID (None to search all folders)
        """
        credential = (
            db.query(GoogleDriveCredential)
            .filter(GoogleDriveCredential.is_active == True)
            .first()
        )

        if not credential:
            raise ValueError("Google Drive not connected")

        credential.folder_id = folder_id
        db.commit()

    @staticmethod
    def search_files(
        db: Session,
        query: str,
        section_type: Optional[str] = None,
        max_results: int = 10,
    ) -> List[GoogleDriveFile]:
        """
        Search Google Drive for files matching the query

        Args:
            db: Database session
            query: Search query
            section_type: Optional section type for context
            max_results: Maximum number of results to return

        Returns:
            List of matching files
        """
        creds = GoogleDriveService._get_credentials(db)
        if not creds:
            raise ValueError("Google Drive not connected")

        # Get folder_id from credentials if set
        credential = (
            db.query(GoogleDriveCredential)
            .filter(GoogleDriveCredential.is_active == True)
            .first()
        )
        folder_id = credential.folder_id if credential else None

        try:
            service = build("drive", "v3", credentials=creds)

            # Build search query
            # Enhance query with section type if provided
            search_query = query
            if section_type:
                search_query = f"{query} {section_type}"

            # Search for documents, PDFs, and presentations
            mime_types = [
                "application/vnd.google-apps.document",
                "application/pdf",
                "application/vnd.google-apps.presentation",
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ]

            mime_query = " or ".join([f"mimeType='{mime}'" for mime in mime_types])
            full_query = f"fullText contains '{search_query}' and ({mime_query})"

            # Add folder constraint if folder_id is set
            if folder_id:
                full_query += f" and '{folder_id}' in parents"

            # Execute search
            results = (
                service.files()
                .list(
                    q=full_query,
                    pageSize=max_results,
                    fields="files(id, name, mimeType, webViewLink, modifiedTime, size, thumbnailLink)",
                    orderBy="modifiedTime desc",
                )
                .execute()
            )

            files = results.get("files", [])

            # Convert to schema objects
            drive_files = []
            for file in files:
                drive_file = GoogleDriveFile(
                    id=file["id"],
                    name=file["name"],
                    mime_type=file["mimeType"],
                    web_view_link=file.get("webViewLink"),
                    modified_time=(
                        datetime.fromisoformat(file["modifiedTime"].replace("Z", "+00:00"))
                        if "modifiedTime" in file
                        else None
                    ),
                    size=file.get("size"),
                    thumbnail_link=file.get("thumbnailLink"),
                )
                drive_files.append(drive_file)

            return drive_files

        except HttpError as error:
            raise ValueError(f"Google Drive API error: {error}")

    @staticmethod
    def get_file_content(db: Session, file_id: str) -> Optional[str]:
        """
        Get the content of a Google Drive file

        Args:
            db: Database session
            file_id: Google Drive file ID

        Returns:
            File content as string, or None if not supported
        """
        creds = GoogleDriveService._get_credentials(db)
        if not creds:
            raise ValueError("Google Drive not connected")

        try:
            service = build("drive", "v3", credentials=creds)

            # Get file metadata
            file = service.files().get(fileId=file_id, fields="mimeType,name").execute()
            mime_type = file.get("mimeType")

            # Handle Google Docs - export to plain text
            if mime_type == "application/vnd.google-apps.document":
                content = (
                    service.files()
                    .export(fileId=file_id, mimeType="text/plain")
                    .execute()
                )
                return content.decode("utf-8")

            # Handle Google Presentations - export to plain text
            if mime_type == "application/vnd.google-apps.presentation":
                content = (
                    service.files()
                    .export(fileId=file_id, mimeType="text/plain")
                    .execute()
                )
                return content.decode("utf-8")

            # Handle PDF files
            if mime_type == "application/pdf":
                request = service.files().get_media(fileId=file_id)
                file_buffer = io.BytesIO()
                downloader = MediaIoBaseDownload(file_buffer, request)
                done = False
                while not done:
                    status, done = downloader.next_chunk()

                file_buffer.seek(0)

                # Extract text using pdfplumber
                try:
                    with pdfplumber.open(file_buffer) as pdf:
                        text = ""
                        for page in pdf.pages:
                            page_text = page.extract_text()
                            if page_text:
                                text += page_text + "\n\n"
                        return text.strip()
                except Exception as e:
                    raise ValueError(f"Error extracting PDF content: {e}")

            # Handle Word documents (.docx)
            if mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                request = service.files().get_media(fileId=file_id)
                file_buffer = io.BytesIO()
                downloader = MediaIoBaseDownload(file_buffer, request)
                done = False
                while not done:
                    status, done = downloader.next_chunk()

                file_buffer.seek(0)

                # Extract text using python-docx
                try:
                    doc = Document(file_buffer)
                    text = ""
                    for paragraph in doc.paragraphs:
                        text += paragraph.text + "\n"
                    return text.strip()
                except Exception as e:
                    raise ValueError(f"Error extracting Word document content: {e}")

            # Handle PowerPoint files (.pptx)
            if mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                request = service.files().get_media(fileId=file_id)
                file_buffer = io.BytesIO()
                downloader = MediaIoBaseDownload(file_buffer, request)
                done = False
                while not done:
                    status, done = downloader.next_chunk()

                file_buffer.seek(0)

                # Extract text from PowerPoint
                try:
                    from pptx import Presentation
                    prs = Presentation(file_buffer)
                    text = ""
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text += shape.text + "\n"
                    return text.strip()
                except ImportError:
                    # python-pptx not installed, return None
                    return None
                except Exception as e:
                    raise ValueError(f"Error extracting PowerPoint content: {e}")

            # Unsupported file type
            return None

        except HttpError as error:
            raise ValueError(f"Google Drive API error: {error}")

    @staticmethod
    def chunk_content(
        content: str,
        chunk_size: int = 1000,
        overlap: int = 200
    ) -> List[Dict[str, Any]]:
        """
        Chunk text content into smaller pieces with overlap

        Args:
            content: Text content to chunk
            chunk_size: Maximum characters per chunk
            overlap: Number of characters to overlap between chunks

        Returns:
            List of chunks with metadata
        """
        if not content:
            return []

        chunks = []
        start = 0
        chunk_index = 0

        while start < len(content):
            # Calculate end position
            end = start + chunk_size

            # Find the last sentence or paragraph boundary before chunk_size
            if end < len(content):
                # Try to break at paragraph
                last_para = content.rfind("\n\n", start, end)
                if last_para > start:
                    end = last_para + 2
                else:
                    # Try to break at sentence
                    last_period = max(
                        content.rfind(". ", start, end),
                        content.rfind(".\n", start, end),
                        content.rfind("? ", start, end),
                        content.rfind("! ", start, end)
                    )
                    if last_period > start:
                        end = last_period + 2

            chunk_text = content[start:end].strip()

            if chunk_text:
                chunks.append({
                    "chunk_index": chunk_index,
                    "text": chunk_text,
                    "start_position": start,
                    "end_position": end,
                    "length": len(chunk_text)
                })
                chunk_index += 1

            # Move start position with overlap
            start = end - overlap if end < len(content) else len(content)

        return chunks
